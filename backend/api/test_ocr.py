from io import BytesIO
from tempfile import NamedTemporaryFile, TemporaryDirectory
from unittest.mock import patch

from django.contrib.auth import get_user_model
from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import SimpleTestCase, TestCase, override_settings
from django.urls import reverse

from .services.document_extractor import extract_document
from .services.ocr_processor import OCRProcessingError, extract_email_addresses, get_reader, read_image


def png_bytes():
    from PIL import Image

    content = BytesIO()
    Image.new("RGB", (40, 20), "white").save(content, format="PNG")
    return content.getvalue()


class EmailExtractionTests(SimpleTestCase):
    def test_addresses_are_validated_normalized_and_deduplicated(self):
        text = (
            "Alıcı: Pilot.Test@Example.COM; kopya: ops@example.org. "
            "Tekrar: pilot.test@example.com ve geçersiz: user@localhost"
        )

        self.assertEqual(
            extract_email_addresses(text),
            ["pilot.test@example.com", "ops@example.org"],
        )

    def test_missing_local_models_raise_setup_error_without_importing_engine(self):
        with TemporaryDirectory() as model_directory, override_settings(
            OCR_MODEL_DIR=model_directory,
            OCR_ALLOW_MODEL_DOWNLOAD=False,
        ):
            get_reader.cache_clear()
            with self.assertRaisesMessage(OCRProcessingError, "model dosyaları bulunamadı"):
                get_reader()
            get_reader.cache_clear()


class DocumentOCRTests(SimpleTestCase):
    @patch("api.services.document_extractor.read_image")
    def test_png_uses_ocr_and_returns_email_metadata(self, mocked_read):
        mocked_read.return_value = "İletişim\npilot@example.com"
        with NamedTemporaryFile(suffix=".png") as image_file:
            image_file.write(png_bytes())
            image_file.flush()

            result = extract_document(image_file.name, use_ocr=True)

        self.assertIn("[OCR - Resim]", result["text"])
        self.assertEqual(result["ocr"]["processed_images"], 1)
        self.assertEqual(result["ocr"]["email_addresses"], ["pilot@example.com"])

    @patch("api.services.document_extractor.read_image")
    def test_multi_frame_tiff_processes_each_frame(self, mocked_read):
        from PIL import Image

        mocked_read.side_effect = ["Birinci", "İkinci"]
        content = BytesIO()
        first = Image.new("RGB", (20, 20), "white")
        second = Image.new("RGB", (20, 20), "black")
        first.save(content, format="TIFF", save_all=True, append_images=[second])
        with NamedTemporaryFile(suffix=".tiff") as image_file:
            image_file.write(content.getvalue())
            image_file.flush()

            result = extract_document(image_file.name, use_ocr=True)

        self.assertEqual(result["ocr"]["processed_images"], 2)
        self.assertIn("[OCR - Kare 2]", result["text"])

    @patch("api.services.document_extractor.read_image")
    def test_docx_embedded_image_is_labeled_and_appended(self, mocked_read):
        from docx import Document

        mocked_read.return_value = "embedded@example.com"
        document = Document()
        document.add_paragraph("Belge metni")
        document.add_picture(BytesIO(png_bytes()))
        with NamedTemporaryFile(suffix=".docx") as document_file:
            document.save(document_file.name)

            result = extract_document(document_file.name, use_ocr=True)

        self.assertIn("Belge metni", result["text"])
        self.assertIn("[OCR - DOCX Görsel 1]", result["text"])
        self.assertEqual(result["ocr"]["email_addresses"], ["embedded@example.com"])

    @patch("api.services.document_extractor.read_image")
    def test_pptx_and_xlsx_images_keep_container_context(self, mocked_read):
        from openpyxl import Workbook
        from openpyxl.drawing.image import Image as WorksheetImage
        from pptx import Presentation
        from pptx.util import Inches

        mocked_read.return_value = "image@example.com"
        with self.subTest("pptx"):
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(BytesIO(png_bytes()), Inches(1), Inches(1))
            with NamedTemporaryFile(suffix=".pptx") as presentation_file:
                presentation.save(presentation_file.name)
                result = extract_document(presentation_file.name, use_ocr=True)
            self.assertIn("[OCR - Slayt 1, Görsel 1]", result["text"])

        with self.subTest("xlsx"):
            workbook = Workbook()
            worksheet = workbook.active
            worksheet.title = "İletişim"
            worksheet.add_image(WorksheetImage(BytesIO(png_bytes())), "A1")
            with NamedTemporaryFile(suffix=".xlsx") as workbook_file:
                workbook.save(workbook_file.name)
                result = extract_document(workbook_file.name, use_ocr=True)
            self.assertIn("[OCR - Çalışma Sayfası İletişim, Görsel 1]", result["text"])

    @patch("api.services.document_extractor.read_image")
    def test_pdf_only_ocrs_pages_without_enough_native_text(self, mocked_read):
        import fitz

        mocked_read.return_value = "scan@example.com"
        pdf = fitz.open()
        text_page = pdf.new_page()
        text_page.insert_text((72, 72), "Bu sayfada OCR gerektirmeyecek kadar uzun yerel metin bulunuyor.")
        pdf.new_page()
        with NamedTemporaryFile(suffix=".pdf") as pdf_file:
            pdf.save(pdf_file.name)
            result = extract_document(pdf_file.name, use_ocr=True)
        pdf.close()

        self.assertEqual(result["ocr"]["processed_pages"], 1)
        self.assertIn("[OCR - Sayfa 2]", result["text"])
        self.assertNotIn("[OCR - Sayfa 1]", result["text"])
        mocked_read.assert_called_once()

    @override_settings(OCR_MAX_PIXELS=10)
    def test_pixel_limit_raises_clear_error(self):
        from PIL import Image

        with self.assertRaisesMessage(OCRProcessingError, "izin verilen sınır"):
            read_image(Image.new("RGB", (10, 10)), "Büyük görsel")


class DocumentUploadOCRApiTests(TestCase):
    def setUp(self):
        self.media_directory = TemporaryDirectory()
        self.settings_override = override_settings(MEDIA_ROOT=self.media_directory.name)
        self.settings_override.enable()
        user = get_user_model().objects.create_user(username="ocr-user", password="StrongPass123!")
        self.client.force_login(user)

    def tearDown(self):
        self.settings_override.disable()
        self.media_directory.cleanup()

    def test_image_requires_ocr(self):
        response = self.client.post(
            reverse("document-upload"),
            data={
                "file": SimpleUploadedFile("mail.png", png_bytes(), content_type="image/png"),
                "prompt": "",
                "use_ocr": "false",
                "use_ai": "false",
            },
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn("use_ocr", response.json())

    def test_ai_requires_prompt(self):
        response = self.client.post(
            reverse("document-upload"),
            data={
                "file": SimpleUploadedFile("not.txt", b"metin", content_type="text/plain"),
                "prompt": "",
                "use_ai": "true",
            },
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn("prompt", response.json())

    @patch("api.views.process_document_text")
    @patch("api.views.extract_document")
    def test_ocr_only_upload_skips_ai_and_returns_metadata(self, mocked_extract, mocked_ai):
        mocked_extract.return_value = {
            "text": "pilot@example.com",
            "ocr": {
                "enabled": True,
                "engine": "easyocr",
                "languages": ["tr", "en"],
                "processed_images": 1,
                "processed_pages": 0,
                "email_addresses": ["pilot@example.com"],
                "warnings": [],
            },
        }

        response = self.client.post(
            reverse("document-upload"),
            data={
                "file": SimpleUploadedFile("mail.png", png_bytes(), content_type="image/png"),
                "prompt": "",
                "use_ocr": "true",
                "use_ai": "false",
            },
        )

        self.assertEqual(response.status_code, 201)
        self.assertFalse(response.json()["ai_result"]["ai_enabled"])
        self.assertEqual(
            response.json()["ai_result"]["ocr"]["email_addresses"],
            ["pilot@example.com"],
        )
        mocked_ai.assert_not_called()

    @patch("api.views.extract_document")
    def test_existing_upload_defaults_remain_ai_on_and_ocr_off(self, mocked_extract):
        mocked_extract.return_value = {
            "text": "mevcut belge metni",
            "ocr": {
                "enabled": False,
                "engine": None,
                "languages": [],
                "processed_images": 0,
                "processed_pages": 0,
                "email_addresses": [],
                "warnings": [],
            },
        }

        response = self.client.post(
            reverse("document-upload"),
            data={
                "file": SimpleUploadedFile("not.txt", b"metin", content_type="text/plain"),
                "prompt": "Özetle",
            },
        )

        self.assertEqual(response.status_code, 201)
        mocked_extract.assert_called_once()
        self.assertEqual(mocked_extract.call_args.kwargs, {"use_ocr": False})
        self.assertTrue(response.json()["ai_result"]["ai_enabled"])
