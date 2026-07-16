from pathlib import Path

from django.conf import settings

from .ocr_processor import (
    OCRProcessingError,
    empty_ocr_metadata,
    extract_email_addresses,
    open_image_bytes,
    read_image,
)


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".csv", ".md", *IMAGE_EXTENSIONS
}


class UnsupportedDocumentError(ValueError):
    pass


def extract_text(file_path):
    return extract_document(file_path, use_ocr=False)["text"]


def extract_document(file_path, use_ocr=False):
    path = Path(file_path)
    extension = path.suffix.lower()

    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise UnsupportedDocumentError(f"Desteklenmeyen dosya tipi: {extension}. Desteklenenler: {supported}")

    metadata = empty_ocr_metadata(enabled=use_ocr)
    if extension in IMAGE_EXTENSIONS:
        if not use_ocr:
            raise UnsupportedDocumentError("Resim dosyalarından metin çıkarmak için OCR etkinleştirilmelidir.")
        text = _extract_image(path, metadata)
        metadata["email_addresses"] = extract_email_addresses(text)
        return {"text": _limit_text(_normalize_text(text)), "ocr": metadata}

    extractors = {
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".pptx": _extract_pptx,
        ".txt": _extract_plain_text,
        ".csv": _extract_plain_text,
        ".md": _extract_plain_text,
    }
    if extension == ".pdf":
        text = _extract_pdf(path, metadata if use_ocr else None)
    else:
        text = extractors[extension](path)
        if use_ocr and extension in {".docx", ".xlsx", ".pptx"}:
            ocr_text = _extract_embedded_images(path, extension, metadata)
            if ocr_text:
                text = f"{text}\n\n{ocr_text}" if text else ocr_text

    metadata["email_addresses"] = extract_email_addresses(text) if use_ocr else []
    return {"text": _limit_text(_normalize_text(text)), "ocr": metadata}


def _extract_pdf(path, metadata=None):
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    page_texts = []
    for index, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        page_texts.append(page_text)
        pages.append(f"[Sayfa {index}]\n{page_text}" if page_text.strip() else "")

    if metadata is not None:
        _append_pdf_ocr(path, page_texts, pages, metadata)
    return "\n\n".join(page for page in pages if page)


def _append_pdf_ocr(path, page_texts, pages, metadata):
    threshold = getattr(settings, "OCR_PDF_MIN_TEXT_LENGTH", 40)
    page_indexes = [index for index, text in enumerate(page_texts) if len(text.strip()) < threshold]
    if not page_indexes:
        return

    try:
        import fitz
    except ImportError as exc:
        raise OCRProcessingError(
            "Taranmış PDF OCR işlemi için PyMuPDF kurulu değil. Bağımlılıkları yeniden kurun."
        ) from exc

    pdf = fitz.open(str(path))
    try:
        dpi = getattr(settings, "OCR_PDF_DPI", 200)
        for page_index in page_indexes:
            if not _has_ocr_capacity(metadata):
                metadata["warnings"].append("OCR öğe sınırına ulaşıldı; kalan PDF sayfaları atlandı.")
                break
            label = f"PDF Sayfa {page_index + 1}"
            try:
                pixmap = pdf.load_page(page_index).get_pixmap(dpi=dpi, alpha=False)
                image = open_image_bytes(pixmap.tobytes("png"))
                ocr_text = read_image(image, label)
                metadata["processed_pages"] += 1
                if ocr_text:
                    ocr_block = f"[OCR - Sayfa {page_index + 1}]\n{ocr_text}"
                    pages[page_index] = (
                        f"{pages[page_index]}\n\n{ocr_block}" if pages[page_index] else ocr_block
                    )
            except Exception as exc:
                metadata["warnings"].append(str(exc))
    finally:
        pdf.close()


def _extract_image(path, metadata):
    from PIL import Image, ImageSequence

    try:
        with Image.open(path) as image:
            blocks = []
            for frame_index, frame in enumerate(ImageSequence.Iterator(image), start=1):
                if not _has_ocr_capacity(metadata):
                    raise OCRProcessingError("OCR öğe sınırına ulaşıldı; resmin tüm kareleri işlenemedi.")
                label = f"Resim {frame_index}"
                text = read_image(frame.copy(), label)
                metadata["processed_images"] += 1
                if text:
                    marker = f"[OCR - Kare {frame_index}]" if getattr(image, "n_frames", 1) > 1 else "[OCR - Resim]"
                    blocks.append(f"{marker}\n{text}")
            return "\n\n".join(blocks)
    except OCRProcessingError:
        raise
    except Exception as exc:
        raise OCRProcessingError(f"Resim OCR ile okunamadı: {exc}") from exc


def _extract_embedded_images(path, extension, metadata):
    iterators = {
        ".docx": _iter_docx_images,
        ".pptx": _iter_pptx_images,
        ".xlsx": _iter_xlsx_images,
    }
    blocks = []
    try:
        for label, content in iterators[extension](path):
            if not _has_ocr_capacity(metadata):
                metadata["warnings"].append("OCR öğe sınırına ulaşıldı; kalan gömülü görseller atlandı.")
                break
            try:
                image = open_image_bytes(content)
                text = read_image(image, label)
                metadata["processed_images"] += 1
                if text:
                    blocks.append(f"[OCR - {label}]\n{text}")
            except Exception as exc:
                metadata["warnings"].append(str(exc))
    except Exception as exc:
        metadata["warnings"].append(f"Gömülü görseller okunamadı: {exc}")
    return "\n\n".join(blocks)


def _iter_docx_images(path):
    from docx import Document as DocxDocument

    document = DocxDocument(str(path))
    image_index = 0
    for relation in document.part.rels.values():
        if "image" not in relation.reltype:
            continue
        image_index += 1
        yield f"DOCX Görsel {image_index}", relation.target_part.blob


def _iter_pptx_images(path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(path))
    for slide_index, slide in enumerate(presentation.slides, start=1):
        image_index = 0
        for shape in _walk_shapes(slide.shapes, MSO_SHAPE_TYPE.GROUP):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            image_index += 1
            yield f"Slayt {slide_index}, Görsel {image_index}", shape.image.blob


def _walk_shapes(shapes, group_shape_type):
    for shape in shapes:
        if shape.shape_type == group_shape_type:
            yield from _walk_shapes(shape.shapes, group_shape_type)
        else:
            yield shape


def _iter_xlsx_images(path):
    from openpyxl import load_workbook

    workbook = load_workbook(filename=str(path), data_only=True, read_only=False)
    try:
        for sheet in workbook.worksheets:
            for image_index, image in enumerate(getattr(sheet, "_images", []), start=1):
                yield f"Çalışma Sayfası {sheet.title}, Görsel {image_index}", image._data()
    finally:
        workbook.close()


def _has_ocr_capacity(metadata):
    processed = metadata["processed_images"] + metadata["processed_pages"]
    return processed < getattr(settings, "OCR_MAX_IMAGES", 50)


def _extract_docx(path):
    from docx import Document as DocxDocument

    document = DocxDocument(str(path))
    blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]

    for table_index, table in enumerate(document.tables, start=1):
        blocks.append(f"[Tablo {table_index}]")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                blocks.append(" | ".join(cells))

    return "\n".join(blocks)


def _extract_xlsx(path):
    from openpyxl import load_workbook

    workbook = load_workbook(filename=str(path), data_only=True, read_only=True)
    blocks = []

    for sheet in workbook.worksheets:
        blocks.append(f"[Sayfa: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                blocks.append(" | ".join(values))

    workbook.close()
    return "\n".join(blocks)


def _extract_pptx(path):
    from pptx import Presentation

    presentation = Presentation(str(path))
    blocks = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        blocks.append(f"[Slayt {slide_index}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                blocks.append(shape.text.strip())

    return "\n".join(blocks)


def _extract_plain_text(path):
    for encoding in ("utf-8", "utf-8-sig", "cp1254", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(errors="ignore")


def _normalize_text(text):
    lines = [line.strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line)


def _limit_text(text):
    max_length = getattr(settings, "DOCUMENT_MAX_TEXT_LENGTH", 120_000)
    if len(text) <= max_length:
        return text
    return text[:max_length] + "\n\n[Metin yerel işlem limiti nedeniyle kısaltıldı.]"
